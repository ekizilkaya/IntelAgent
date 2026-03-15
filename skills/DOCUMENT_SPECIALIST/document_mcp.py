import json
import os
import sys
from mcp.server.fastmcp import FastMCP

# Initialize FastMCP Server
mcp = FastMCP("DOCUMENT_SPECIALIST")

@mcp.tool()
def create_excel(filepath: str, data_json: str, sheet_name: str = "Sheet1") -> str:
    """
    Creates an Excel spreadsheet from JSON data.
    
    Args:
        filepath: Absolute path to save the .xlsx file.
        data_json: A JSON string representing a list of dictionaries (rows).
        sheet_name: Name of the sheet.
    """
    try:
        import pandas as pd
        data = json.loads(data_json)
        df = pd.DataFrame(data)
        
        # Ensure the directory exists
        os.makedirs(os.path.dirname(os.path.abspath(filepath)), exist_ok=True)
        
        # Use ExcelWriter with Engine openpyxl to auto-adjust columns (optional but nice)
        with pd.ExcelWriter(filepath, engine='openpyxl') as writer:
            df.to_excel(writer, sheet_name=sheet_name, index=False)
            
        return f"Successfully created Excel file at: {filepath}"
    except ImportError:
        return "Error: pandas or openpyxl missing. Please install: pip install pandas openpyxl"
    except Exception as e:
        return f"Error creating Excel file: {e}"

@mcp.tool()
def read_excel_or_csv(filepath: str, max_rows: int = 50) -> str:
    """
    Reads data from an Excel (.xlsx, .xls) or CSV (.csv) file and returns a JSON string,
    or a generic markdown summary. Useful for summarizing uploaded datasets or spreadsheets.
    
    Args:
        filepath: Absolute path to the .xlsx or .csv file.
        max_rows: Maximum number of rows to return to prevent context overflow.
    """
    try:
        import pandas as pd
    except ImportError:
        return "Error: pandas missing. pip install pandas"
        
    try:
        if not os.path.exists(filepath):
            return f"Error: File not found at {filepath}"
            
        if filepath.endswith((".csv", ".txt")):
            df = pd.read_csv(filepath)
        else:
            # Using openpyxl for xlsx and ignoring data it can't parse like charts
            df = pd.read_excel(filepath, engine='openpyxl')
            
        total_rows, total_cols = df.shape
        preview_df = df.head(max_rows)
        
        # Convert all to string to avoid serialization issues with dates/complex objects
        preview_df = preview_df.astype(str)
        data_json = preview_df.to_json(orient="records", force_ascii=False)
        
        summary = (
            f"File: {os.path.basename(filepath)}\n"
            f"Total Rows: {total_rows}\n"
            f"Total Columns: {total_cols}\n"
            f"Columns: {', '.join(df.columns.astype(str).tolist())}\n\n"
            f"Data Preview (First {max_rows} rows):\n{data_json}"
        )
        return summary
    except Exception as e:
        import traceback
        return f"Error reading Excel/CSV file details: {e}\nTraceback: {traceback.format_exc()}"

@mcp.tool()
def create_powerpoint(filepath: str, slides_json: str) -> str:
    """
    Creates a PowerPoint presentation with custom typography.
    
    Args:
        filepath: Absolute path to save the .pptx file.
        slides_json: JSON string of a list of slide objects, e.g., [{"title": "Intro", "content": "Hello", "font_name": "Arial", "font_size": 24}]
    """
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError:
        return "Error: python-pptx missing. Please install: pip install python-pptx"
        
    try:
        data = json.loads(slides_json)
        prs = Presentation()
        
        # Ensure the directory exists
        os.makedirs(os.path.dirname(os.path.abspath(filepath)), exist_ok=True)

        for slide_data in data:
            title_text = slide_data.get("title", "")
            content_text = slide_data.get("content", "")
            font_name = slide_data.get("font_name", "Calibri")
            font_size = slide_data.get("font_size", 18)
            
            # Use blank layout with title & content
            title_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(title_slide_layout)
            
            # Title
            title = slide.shapes.title
            if title:
                title.text = title_text
            
            # Content
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.text = content_text
            
            # Apply styling
            for p in tf.paragraphs:
                for r in p.runs:
                    r.font.name = font_name
                    r.font.size = Pt(font_size)

        prs.save(filepath)
        return f"Successfully created PowerPoint presentation at: {filepath}"
    except Exception as e:
        return f"Error creating PowerPoint: {e}"

@mcp.tool()
def extract_pdf_text(filepath: str) -> str:
    """
    Extracts text from a given PDF file.
    
    Args:
        filepath: Absolute path to the .pdf file.
    """
    try:
        import fitz  # PyMuPDF
    except ImportError:
        return "Error: PyMuPDF missing. Please install: pip install PyMuPDF"
        
    try:
        doc = fitz.open(filepath)
        text = []
        for i, page in enumerate(doc):
            text.append(f"--- Page {i+1} ---\n{page.get_text()}")
        return "\n".join(text)
    except Exception as e:
        return f"Error reading PDF {filepath}: {e}"

@mcp.tool()
def merge_pdfs(filepaths_json: str, output_filepath: str) -> str:
    """
    Merges multiple PDF files into one.
    
    Args:
        filepaths_json: JSON string array of absolute PDF file paths.
        output_filepath: Absolute path to save the merged .pdf.
    """
    try:
        import fitz  # PyMuPDF
    except ImportError:
        return "Error: PyMuPDF missing. Please install: pip install PyMuPDF"
        
    try:
        paths = json.loads(filepaths_json)
        merged_doc = fitz.open()
        
        # Ensure the directory exists
        os.makedirs(os.path.dirname(os.path.abspath(output_filepath)), exist_ok=True)
        
        for path in paths:
            doc = fitz.open(path)
            merged_doc.insert_pdf(doc)
            
        merged_doc.save(output_filepath)
        return f"Successfully merged PDFs into: {output_filepath}"
    except Exception as e:
        return f"Error merging PDFs: {e}"

@mcp.tool()
def create_word_document(filepath: str, document_json: str) -> str:
    """
    Creates a professional Word (.docx) document with proper formatting.
    Ideal for business proposals, NDAs, HR reviews, meeting minutes, and job descriptions.
    
    Args:
        filepath: Absolute path to save the .docx file.
        document_json: JSON string representing the document structure. Should be a list of block objects.
                       Supported blocks:
                       - {"type": "heading", "level": 1|2|3, "text": "..."}
                       - {"type": "paragraph", "text": "..."}
                       - {"type": "bullet_point", "text": "..."}
                       - {"type": "numbered_point", "text": "..."}
                       - {"type": "bold_paragraph", "text": "..."}
    """
    try:
        import docx
        from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
    except ImportError:
        return "Error: python-docx missing. Please install: pip install python-docx"
        
    try:
        data = json.loads(document_json)
        doc = docx.Document()
        
        os.makedirs(os.path.dirname(os.path.abspath(filepath)), exist_ok=True)
        
        for block in data:
            b_type = block.get("type", "paragraph")
            text = block.get("text", "")
            
            if b_type == "heading":
                level = block.get("level", 1)
                doc.add_heading(text, level=level)
            elif b_type == "bullet_point":
                doc.add_paragraph(text, style="List Bullet")
            elif b_type == "numbered_point":
                doc.add_paragraph(text, style="List Number")
            elif b_type == "bold_paragraph":
                p = doc.add_paragraph()
                p.add_run(text).bold = True
            else: # paragraph default
                doc.add_paragraph(text)
                
        doc.save(filepath)
        return f"Successfully created Word document at: {filepath}"
    except Exception as e:
        return f"Error creating Word document: {e}"

if __name__ == "__main__":
    # Start the FastMCP stdio server
    mcp.run()
